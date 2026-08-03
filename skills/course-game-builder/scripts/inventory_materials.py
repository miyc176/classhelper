#!/usr/bin/env python3
"""Inventory course materials and extract machine-readable text plus visual review tasks."""

from __future__ import annotations

import argparse
import concurrent.futures
import hashlib
import json
import re
import shutil
import sys
import time
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET


SUPPORTED = {".pptx", ".pdf", ".docx", ".txt", ".md", ".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}


@dataclass
class ExtractedSource:
    source_id: str
    path: str
    kind: str
    extraction_status: str = "complete"
    notes: str = ""
    page_or_slide_count: int = 0
    text_units: list[dict[str, Any]] = field(default_factory=list)
    visual_units: list[dict[str, Any]] = field(default_factory=list)
    context_units: list[dict[str, Any]] = field(default_factory=list)


def stable_id(index: int) -> str:
    return f"src_{index:03d}"


def clean_text(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip()


def sha1(path: Path) -> str:
    h = hashlib.sha1()
    with path.open("rb") as handle:
      for chunk in iter(lambda: handle.read(1024 * 1024), b""):
          h.update(chunk)
    return h.hexdigest()


def write_visual_copy(source_id: str, locator: str, image_name: str, data: bytes, out_dir: Path) -> str:
    safe_locator = re.sub(r"[^a-zA-Z0-9_-]+", "-", locator).strip("-").lower() or "visual"
    target_dir = out_dir / "visuals" / source_id
    target_dir.mkdir(parents=True, exist_ok=True)
    suffix = Path(image_name).suffix or ".bin"
    target = target_dir / f"{safe_locator}-{len(list(target_dir.iterdir())) + 1:03d}{suffix}"
    target.write_bytes(data)
    return str(target)


def write_context_visual(source_id: str, unit_index: int, visual_index: int, suffix: str, data: bytes, out_dir: Path) -> str:
    target_dir = out_dir / "visuals" / source_id
    target_dir.mkdir(parents=True, exist_ok=True)
    target = target_dir / f"unit-{unit_index:04d}-visual-{visual_index:03d}.{suffix.lstrip('.') or 'bin'}"
    if not target.exists() or target.read_bytes() != data:
        target.write_bytes(data)
    return str(target)


def normalized_box(shape: Any, width: int, height: int) -> list[float]:
    return [
        round(float(shape.left) / width, 5),
        round(float(shape.top) / height, 5),
        round(float(shape.width) / width, 5),
        round(float(shape.height) / height, 5),
    ]


def box_distance(first: list[float], second: list[float]) -> float:
    first_x, first_y = first[0] + first[2] / 2, first[1] + first[3] / 2
    second_x, second_y = second[0] + second[2] / 2, second[1] + second[3] / 2
    return (first_x - second_x) ** 2 + (first_y - second_y) ** 2


def xml_text(xml_bytes: bytes) -> str:
    try:
        root = ET.fromstring(xml_bytes)
    except ET.ParseError:
        return ""
    texts = [node.text for node in root.iter() if node.tag.endswith("}t") and node.text]
    return clean_text(" ".join(texts))


def extract_pptx(path: Path, source_id: str, out_dir: Path) -> ExtractedSource:
    source = ExtractedSource(source_id, str(path), "pptx")
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(str(path))
        source.page_or_slide_count = len(prs.slides)
        image_assets: dict[str, dict[str, str]] = {}
        for idx, slide in enumerate(prs.slides, start=1):
            locator = f"slide {idx}"
            parts: list[str] = []
            elements: list[dict[str, Any]] = []
            picture_elements: list[dict[str, Any]] = []
            text_elements: list[dict[str, Any]] = []
            connector_elements: list[dict[str, Any]] = []

            def visit(shapes: Any, parent_id: str | None = None) -> None:
                for z_index, shape in enumerate(shapes):
                    element_id = f"s{idx:04d}-e{len(elements) + 1:03d}"
                    shape_type = str(getattr(shape, "shape_type", "unknown"))
                    box = normalized_box(shape, prs.slide_width, prs.slide_height)
                    element: dict[str, Any] = {
                        "id": element_id,
                        "shape_type": shape_type,
                        "name": str(getattr(shape, "name", "")),
                        "bbox_norm": box,
                        "z_index": z_index,
                    }
                    try:
                        alt_text = str(shape._element.xpath(".//p:cNvPr")[0].get("descr") or "").strip()
                    except Exception:
                        alt_text = ""
                    if alt_text:
                        element["alt_text"] = alt_text
                    if parent_id:
                        element["parent_group_id"] = parent_id
                    text = clean_text(str(getattr(shape, "text", "") or ""))
                    if getattr(shape, "has_table", False):
                        table_rows = [" | ".join(clean_text(cell.text) for cell in row.cells) for row in shape.table.rows]
                        element["table_rows"] = table_rows
                        text = clean_text(" ".join([text, *table_rows]))
                    if getattr(shape, "has_chart", False):
                        chart_series = []
                        for series in shape.chart.series:
                            chart_series.append({
                                "name": str(getattr(series, "name", "")),
                                "values": [value for value in getattr(series, "values", [])],
                            })
                        element["chart_series"] = chart_series
                    if text:
                        element["text"] = text
                        parts.append(text)
                        text_elements.append(element)
                    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                        visual_index = len(picture_elements) + 1
                        image = shape.image
                        asset_sha1 = hashlib.sha1(image.blob).hexdigest()
                        existing_asset = image_assets.get(asset_sha1)
                        extracted_path = existing_asset["extracted_path"] if existing_asset else write_context_visual(source_id, idx, visual_index, image.ext, image.blob, out_dir)
                        element.update({"modality": "embedded_image", "extracted_path": extracted_path, "asset_sha1": asset_sha1})
                        if existing_asset:
                            element["duplicate_of"] = existing_asset["element_id"]
                        else:
                            image_assets[asset_sha1] = {"extracted_path": extracted_path, "element_id": element_id}
                        picture_elements.append(element)
                    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.LINE:
                        connector_elements.append(element)
                    elements.append(element)
                    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                        visit(shape.shapes, element_id)

            visit(slide.shapes)
            relations = []
            for visual_index, picture in enumerate(picture_elements, start=1):
                nearby = sorted(text_elements, key=lambda item: box_distance(picture["bbox_norm"], item["bbox_norm"]))[:3]
                picture["nearby_text_ids"] = [item["id"] for item in nearby]
                picture["nearby_text"] = [item["text"] for item in nearby]
                relations.extend({"type": "spatial_near", "from": picture["id"], "to": item["id"]} for item in nearby)
                source.visual_units.append({
                    "locator": f"{locator} image {visual_index}",
                    "parent_unit": locator,
                    "element_id": picture["id"],
                    "modality": "embedded_image",
                    "bbox_norm": picture["bbox_norm"],
                    "nearby_text": picture["nearby_text"],
                    "extracted_path": picture["extracted_path"],
                    "asset_sha1": picture["asset_sha1"],
                    "duplicate_of": picture.get("duplicate_of", ""),
                    "review_status": "review_in_page_context",
                })
            visual_pairs: set[tuple[str, str]] = set()
            for picture in picture_elements:
                neighbors = sorted(
                    (item for item in picture_elements if item["id"] != picture["id"]),
                    key=lambda item: box_distance(picture["bbox_norm"], item["bbox_norm"]),
                )[:2]
                for neighbor in neighbors:
                    pair = tuple(sorted([picture["id"], neighbor["id"]]))
                    if pair in visual_pairs:
                        continue
                    visual_pairs.add(pair)
                    first_box, second_box = picture["bbox_norm"], neighbor["bbox_norm"]
                    dx = abs((first_box[0] + first_box[2] / 2) - (second_box[0] + second_box[2] / 2))
                    dy = abs((first_box[1] + first_box[3] / 2) - (second_box[1] + second_box[3] / 2))
                    relations.append({
                        "type": "visual_proximity_candidate",
                        "from": pair[0],
                        "to": pair[1],
                        "layout": "side_by_side" if dx >= dy else "stacked",
                        "requires_render_confirmation": True,
                    })
            for connector in connector_elements:
                candidates = [item for item in elements if item["id"] != connector["id"] and item not in connector_elements]
                endpoints = sorted(candidates, key=lambda item: box_distance(connector["bbox_norm"], item["bbox_norm"]))[:2]
                if len(endpoints) == 2:
                    relations.append({"type": "connector_near", "from": endpoints[0]["id"], "to": endpoints[1]["id"], "via": connector["id"]})
            notes = ""
            try:
                notes = clean_text(slide.notes_slide.notes_text_frame.text)
            except Exception:
                notes = ""
            if notes:
                parts.append(notes)
            source.context_units.append({
                "locator": locator,
                "modality": "slide_context",
                "text": clean_text(" ".join(parts)),
                "notes": notes,
                "elements": elements,
                "relations": relations,
                "visual_count": len(picture_elements),
                "review_status": "needs_page_context_inspection",
            })
            text = clean_text(" ".join(parts))
            if text:
                source.text_units.append({
                    "locator": locator,
                    "modality": "text",
                    "text": text,
                })
    except Exception as exc:
        source.extraction_status = "partial"
        source.notes = f"python-pptx extraction failed: {exc}"
    return source


def extract_docx(path: Path, source_id: str, out_dir: Path) -> ExtractedSource:
    source = ExtractedSource(source_id, str(path), "docx")
    try:
        from docx import Document

        doc = Document(str(path))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        text = clean_text(" ".join(parts))
        if text:
            source.text_units.append({"locator": "document body", "modality": "text", "text": text})
        source.context_units.append({
            "locator": "document body", "modality": "document_context", "text": text,
            "review_status": "needs_document_context_inspection",
        })
    except Exception as exc:
        source.extraction_status = "partial"
        source.notes = f"python-docx extraction failed: {exc}"

    try:
        with zipfile.ZipFile(path) as zf:
            for name in sorted(zf.namelist()):
                if name.startswith("word/media/"):
                    copied = write_visual_copy(source_id, "embedded", Path(name).name, zf.read(name), out_dir)
                    source.visual_units.append({
                        "locator": Path(name).name,
                        "parent_unit": "document body",
                        "modality": "embedded_image",
                        "extracted_path": copied,
                        "review_status": "needs_visual_inspection",
                    })
    except Exception as exc:
        source.extraction_status = "partial"
        source.notes = (source.notes + "; " if source.notes else "") + f"docx package scan failed: {exc}"
    return source


def extract_pdf(path: Path, source_id: str, out_dir: Path) -> ExtractedSource:
    source = ExtractedSource(source_id, str(path), "pdf")
    try:
        import pdfplumber

        with pdfplumber.open(str(path)) as pdf:
            source.page_or_slide_count = len(pdf.pages)
            for idx, page in enumerate(pdf.pages, start=1):
                locator = f"page {idx}"
                text = clean_text(page.extract_text() or "")
                if text:
                    source.text_units.append({"locator": locator, "modality": "text", "text": text})
                page_elements = []
                for image_idx, image in enumerate(page.images, start=1):
                    bbox = [image.get("x0"), image.get("top"), image.get("x1"), image.get("bottom")]
                    page_elements.append({"id": f"p{idx:04d}-image-{image_idx:03d}", "modality": "pdf_image_or_region", "bbox": bbox})
                    source.visual_units.append({
                        "locator": f"page {idx} image {image_idx}",
                        "parent_unit": locator,
                        "modality": "pdf_image_or_region",
                        "bbox": bbox,
                        "review_status": "review_in_page_context",
                    })
                source.context_units.append({
                    "locator": locator,
                    "modality": "page_context",
                    "text": text,
                    "elements": page_elements,
                    "visual_count": len(page_elements),
                    "review_status": "needs_page_context_inspection",
                })
    except Exception as exc:
        source.extraction_status = "partial"
        source.notes = f"pdfplumber extraction failed: {exc}"
    return source


def extract_plain(path: Path, source_id: str) -> ExtractedSource:
    source = ExtractedSource(source_id, str(path), path.suffix.lower().lstrip("."))
    try:
        text = clean_text(path.read_text(encoding="utf-8", errors="ignore"))
        if text:
            source.text_units.append({"locator": "file", "modality": "text", "text": text})
    except Exception as exc:
        source.extraction_status = "blocked"
        source.notes = f"text read failed: {exc}"
    return source


def extract_image(path: Path, source_id: str, out_dir: Path) -> ExtractedSource:
    source = ExtractedSource(source_id, str(path), path.suffix.lower().lstrip("."))
    target_dir = out_dir / "visuals" / source_id
    target_dir.mkdir(parents=True, exist_ok=True)
    target = target_dir / path.name
    if path.resolve() != target.resolve():
        shutil.copy2(path, target)
    source.visual_units.append({
        "locator": "image",
        "parent_unit": "image",
        "modality": "standalone_image",
        "extracted_path": str(target),
        "review_status": "needs_visual_inspection",
    })
    source.context_units.append({
        "locator": "image", "modality": "standalone_image_context", "text": "",
        "elements": [{"id": "image-001", "modality": "standalone_image", "extracted_path": str(target)}],
        "visual_count": 1, "review_status": "needs_page_context_inspection",
    })
    source.page_or_slide_count = 1
    return source


def extract_file(path: Path, source_id: str, out_dir: Path) -> ExtractedSource:
    ext = path.suffix.lower()
    if ext == ".pptx":
        return extract_pptx(path, source_id, out_dir)
    if ext == ".docx":
        return extract_docx(path, source_id, out_dir)
    if ext == ".pdf":
        return extract_pdf(path, source_id, out_dir)
    if ext in {".txt", ".md"}:
        return extract_plain(path, source_id)
    if ext in IMAGE_EXTS:
        return extract_image(path, source_id, out_dir)
    return ExtractedSource(source_id, str(path), ext.lstrip(".") or "other", "blocked", "Unsupported file type")


def discover(input_path: Path) -> list[Path]:
    if input_path.is_file():
        return [input_path]
    return sorted(path for path in input_path.rglob("*") if path.is_file() and path.suffix.lower() in SUPPORTED)


def cached_source(previous: dict[str, Any], path: Path, source_id: str, digest: str) -> ExtractedSource | None:
    if previous.get("extractor_version") != "context-v4":
        return None
    inventory = previous.get("source_inventory") or []
    match = next((item for item in inventory if str(item.get("path")) == str(path) and item.get("sha1") == digest), None)
    if not match:
        return None
    old_id = str(match.get("source_id"))
    visual_units = [{key: value for key, value in item.items() if key != "source_id"} for item in previous.get("visual_units") or [] if str(item.get("source_id")) == old_id]
    if any(item.get("extracted_path") and not Path(item["extracted_path"]).is_file() for item in visual_units):
        return None
    return ExtractedSource(
        source_id=source_id,
        path=str(path),
        kind=str(match.get("kind", path.suffix.lstrip("."))),
        extraction_status=str(match.get("extraction_status", "complete")),
        notes=str(match.get("notes", "")),
        page_or_slide_count=int(match.get("page_or_slide_count", 0)),
        text_units=[{key: value for key, value in item.items() if key != "source_id"} for item in previous.get("text_units") or [] if str(item.get("source_id")) == old_id],
        visual_units=visual_units,
        context_units=[{key: value for key, value in item.items() if key != "source_id"} for item in previous.get("context_units") or [] if str(item.get("source_id")) == old_id],
    )


def build_report(files: list[Path], out_dir: Path, previous: dict[str, Any] | None = None, workers: int = 4) -> dict[str, Any]:
    previous = previous or {}
    started = time.perf_counter()
    file_hashes = {str(path): sha1(path) for path in files}

    def process(item: tuple[int, Path]) -> tuple[int, ExtractedSource, bool]:
        index, path = item
        source_id = stable_id(index)
        cached = cached_source(previous, path, source_id, file_hashes[str(path)])
        return index, cached or extract_file(path, source_id, out_dir), cached is not None

    indexed_files = list(enumerate(files, start=1))
    if workers > 1 and len(indexed_files) > 1:
        with concurrent.futures.ThreadPoolExecutor(max_workers=min(workers, len(indexed_files))) as executor:
            processed = list(executor.map(process, indexed_files))
    else:
        processed = [process(item) for item in indexed_files]
    processed.sort(key=lambda item: item[0])
    sources = [item[1] for item in processed]
    cache_hits = sum(item[2] for item in processed)
    inventory = []
    text_units = []
    visual_units = []
    context_units = []
    coverage_audit = []

    for source in sources:
        inventory.append({
            "source_id": source.source_id,
            "path": source.path,
            "kind": source.kind,
            "sha1": file_hashes[source.path],
            "page_or_slide_count": source.page_or_slide_count,
            "extraction_status": source.extraction_status,
            "notes": source.notes,
        })
        for unit in source.text_units:
            text_units.append({"source_id": source.source_id, **unit})
        for unit in source.visual_units:
            visual_units.append({"source_id": source.source_id, **unit})
        for unit in source.context_units:
            context_units.append({"source_id": source.source_id, **unit})

        locators = {unit["locator"] for unit in source.context_units}
        if not locators:
            locators = {unit["locator"] for unit in source.text_units}
            locators.update(unit.get("parent_unit") or unit["locator"] for unit in source.visual_units)
        if source.kind == "pptx":
            locators.update(f"slide {index}" for index in range(1, source.page_or_slide_count + 1))
        elif source.kind == "pdf":
            locators.update(f"page {index}" for index in range(1, source.page_or_slide_count + 1))
        if not locators:
            coverage_audit.append({
                "source_id": source.source_id,
                "unit": "file",
                "status": "blocked" if source.extraction_status == "blocked" else "no_machine_readable_content",
                "knowledge_ids": [],
                "notes": source.notes or "No text or visual unit was extracted.",
            })
        else:
            for locator in sorted(locators):
                has_visual = any((unit.get("parent_unit") or unit["locator"]) == locator for unit in source.visual_units)
                coverage_audit.append({
                    "source_id": source.source_id,
                    "unit": locator,
                    "status": "needs_visual_review" if has_visual else "needs_knowledge_extraction",
                    "knowledge_ids": [],
                    "notes": "Inspect the full page/slide context before using individual visuals." if has_visual else "",
                })

    visual_asset_hashes = [str(item.get("asset_sha1")) for item in visual_units if item.get("asset_sha1")]
    return {
        "schema_version": "2.0",
        "extractor_version": "context-v4",
        "source_inventory": inventory,
        "text_units": text_units,
        "visual_units": visual_units,
        "context_units": context_units,
        "coverage_audit": coverage_audit,
        "performance": {
            "elapsed_seconds": round(time.perf_counter() - started, 3),
            "workers": workers,
            "cache_hits": cache_hits,
            "cache_misses": len(files) - cache_hits,
            "visual_occurrences": len(visual_units),
            "unique_visual_assets": len(set(visual_asset_hashes)) if visual_asset_hashes else len(visual_units),
            "duplicate_visual_occurrences": len(visual_asset_hashes) - len(set(visual_asset_hashes)),
        },
        "agent_next_steps": [
            "Inspect context_units and rendered whole pages/slides first; use visual_units only for local zoom.",
            "Preserve spatial, connector, group, nearby-text, and cross-visual relationships in knowledge points.",
            "Convert text_units and contextual visual observations into references/knowledge-schema.md knowledge_points.",
            "Replace coverage_audit statuses with covered/no_instructional_content/blocked before game generation.",
        ],
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Inventory course materials for course-game-builder.")
    parser.add_argument("input", help="File or directory containing course materials.")
    parser.add_argument("--out", required=True, help="Output directory for extraction artifacts.")
    parser.add_argument("--workers", type=int, default=4, help="Number of files to extract concurrently.")
    parser.add_argument("--no-cache", action="store_true", help="Ignore a previous matching extraction manifest.")
    args = parser.parse_args()

    input_path = Path(args.input).resolve()
    out_dir = Path(args.out).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    report_path = out_dir / "material-extraction.json"
    previous = {}
    if report_path.is_file() and not args.no_cache:
        try:
            previous = json.loads(report_path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            previous = {}
    files = discover(input_path)
    if not files:
        print(json.dumps({"status": "fail", "error": f"No supported files found in {input_path}"}, ensure_ascii=False, indent=2))
        return 1

    report = build_report(files, out_dir, previous, max(1, args.workers))
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({
        "status": "pass", "report": str(report_path), "sources": len(report["source_inventory"]),
        "performance": report["performance"],
    }, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    sys.exit(main())
